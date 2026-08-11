from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


OUTPUT = Path("C:/repos/auricrux-app/FCA_Auricrux_Zacua_Deck_v4_2.pptx")
BADGE = Path("C:/repos/auricrux-app/apps/auricrux-mobile/assets/auricrux-badge.png")

NAVY = RGBColor(10, 24, 51)
BLUE = RGBColor(20, 95, 170)
TEAL = RGBColor(12, 145, 130)
PURPLE = RGBColor(120, 90, 185)
WHITE = RGBColor(255, 255, 255)
BG = RGBColor(246, 249, 253)
TEXT = RGBColor(28, 34, 44)
MUTED = RGBColor(94, 107, 127)


def style_run(run, size=18, bold=False, color=TEXT):
    run.font.name = "Aptos"
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color


def add_bg(slide, dark=False):
    rect = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0), Inches(0), Inches(13.33), Inches(7.5))
    rect.fill.solid()
    rect.fill.fore_color.rgb = NAVY if dark else BG
    rect.line.fill.background()

    top = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0), Inches(0), Inches(13.33), Inches(0.24))
    top.fill.solid()
    top.fill.fore_color.rgb = TEAL if dark else NAVY
    top.line.fill.background()


def title_block(slide, title, subtitle=None, dark=False):
    tf = slide.shapes.add_textbox(Inches(0.78), Inches(0.42), Inches(11.8), Inches(1.4)).text_frame
    tf.clear()

    p1 = tf.paragraphs[0]
    r1 = p1.add_run()
    r1.text = title
    style_run(r1, size=35, bold=True, color=WHITE if dark else NAVY)

    if subtitle:
        p2 = tf.add_paragraph()
        p2.space_before = Pt(6)
        r2 = p2.add_run()
        r2.text = subtitle
        style_run(r2, size=15, color=RGBColor(203, 221, 242) if dark else MUTED)


def bullets(slide, items, x=0.9, y=1.9, w=11.6, h=4.9, size=22, dark=False):
    tf = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h)).text_frame
    tf.clear()
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.space_after = Pt(10)
        p.font.name = "Aptos"
        p.font.size = Pt(size)
        p.font.color.rgb = WHITE if dark else TEXT


def footer(slide, left, right=""):
    l = slide.shapes.add_textbox(Inches(0.78), Inches(7.02), Inches(8.8), Inches(0.3)).text_frame
    l.clear()
    lp = l.paragraphs[0]
    lr = lp.add_run()
    lr.text = left
    style_run(lr, size=10, color=BLUE)

    r = slide.shapes.add_textbox(Inches(9.3), Inches(7.02), Inches(3.2), Inches(0.3)).text_frame
    r.clear()
    rp = r.paragraphs[0]
    rp.alignment = PP_ALIGN.RIGHT
    rr = rp.add_run()
    rr.text = right
    style_run(rr, size=10, color=MUTED)


def card(slide, x, y, w, h, title, value, body, accent):
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = WHITE
    shape.line.color.rgb = accent
    shape.line.width = Pt(1.5)

    tf = slide.shapes.add_textbox(Inches(x + 0.2), Inches(y + 0.15), Inches(w - 0.4), Inches(h - 0.3)).text_frame
    tf.clear()

    p0 = tf.paragraphs[0]
    r0 = p0.add_run()
    r0.text = title
    style_run(r0, size=12, bold=True, color=accent)

    p1 = tf.add_paragraph()
    p1.space_before = Pt(5)
    r1 = p1.add_run()
    r1.text = value
    style_run(r1, size=25, bold=True, color=NAVY)

    p2 = tf.add_paragraph()
    p2.space_before = Pt(5)
    r2 = p2.add_run()
    r2.text = body
    style_run(r2, size=11, color=MUTED)


def new_slide(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s)
    return s


def build() -> Path:
    prs = Presentation()

    # Cover
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s, dark=True)
    title_block(s, "Future Contractors of America", "Investor Deck v4.2 | Zacua-focused narrative", dark=True)

    if BADGE.exists():
        s.shapes.add_picture(str(BADGE), Inches(9.45), Inches(1.0), height=Inches(1.8))

    bullets(
        s,
        [
            "AI-powered construction operating system for contractors.",
            "FCA unifies operations; Auricrux guides teams on what to do next in context.",
        ],
        y=2.5,
        h=2.4,
        size=23,
        dark=True,
    )
    footer(s, "www.futurecontractorsofamerica.com", "v4.2")

    # Executive summary
    s = new_slide(prs)
    title_block(s, "Executive Summary", "A clear path from fragmented workflows to unified execution")
    bullets(
        s,
        [
            "Problem: contractors still run core operations across disconnected tools and manual follow-up.",
            "Solution: one operating system plus embedded AI guidance for day-to-day execution decisions.",
            "Why now: rising cost pressure and AI readiness create a timing advantage for a category leader.",
        ],
    )
    footer(s, "Category thesis: execution continuity as the defensible wedge", "01")

    # Pain and ROI
    s = new_slide(prs)
    title_block(s, "Pain to ROI", "FCA is positioned around measurable operating outcomes")
    card(s, 0.9, 1.95, 3.95, 2.05, "Leakage Reduction", "18-32%", "Live ROI framing in executive pathway", BLUE)
    card(s, 5.0, 1.95, 3.95, 2.05, "Admin Friction", "$29,825/mo", "Example monthly operational friction", TEAL)
    card(s, 9.1, 1.95, 3.25, 2.05, "Projected Reclaim", "$10,141/mo", "Example monthly value reclaim", PURPLE)

    bullets(
        s,
        [
            "Value is framed in business terms: less leakage, better coordination, stronger billing confidence.",
            "Auricrux shifts teams from reactive status chasing to guided execution with clear ownership.",
        ],
        y=4.45,
        h=2.0,
        size=20,
    )
    footer(s, "Source: futurecontractorsofamerica.com (public experience snapshot)", "02")

    # Product depth
    s = new_slide(prs)
    title_block(s, "Product Scope", "Integrated command layer across the contractor lifecycle")
    bullets(
        s,
        [
            "Qualification, estimating, project flow, field tasks, scheduling, plans/design, and review workflows.",
            "Finance, billing, legal, customer portal, and trust/governance controls in the same product shell.",
            "Academy and workforce pathways linked to live execution, not separate training silos.",
        ],
    )
    footer(s, "Outcome: one operational surface instead of disconnected point tools", "03")

    # Defensibility
    s = new_slide(prs)
    title_block(s, "Defensibility", "Workflow depth + embedded intelligence + continuity")
    bullets(
        s,
        [
            "Workflow depth across preconstruction, delivery, and post-handover continuity.",
            "Auricrux guidance inside active workflows, where teams already do daily work.",
            "A growing operating context that improves recommendation quality over time.",
        ],
    )
    footer(s, "Defensibility model: system-of-action + compounding workflow data", "04")

    # Execution evidence
    s = new_slide(prs)
    title_block(s, "Execution Evidence", "Public founder build velocity and shipping cadence")
    card(s, 0.9, 1.95, 3.95, 1.9, "Contributions", "1,056", "Public GitHub contributions (last year)", BLUE)
    card(s, 5.0, 1.95, 3.95, 1.9, "Recent Commits", "60 in July", "Across five repositories", TEAL)
    card(s, 9.1, 1.95, 3.25, 1.9, "PR Throughput", "33 merged", "In a highly active primary lane", PURPLE)

    bullets(
        s,
        [
            "This is active product execution with sustained iteration, not pre-launch concept work.",
            "Convergence direction centers on FCA Ecosystem as the unified universal repository.",
        ],
        y=4.25,
        h=2.2,
        size=20,
    )
    footer(s, "Source: github.com/Auricrux + FCA org snapshot", "05")

    # GTM
    s = new_slide(prs)
    title_block(s, "Go-To-Market", "Land where continuity failures are most expensive")
    bullets(
        s,
        [
            "Wedge: bid-to-project transition, document coordination, and communication integrity.",
            "Expansion: finance readiness, customer transparency, compliance, and workforce enablement.",
            "Retention: become the daily execution command layer for office and field teams.",
            "Business model: subscription by team and workflow depth, with module expansion over time.",
        ],
    )
    footer(s, "Adoption model: operational pain -> measurable ROI -> platform expansion", "06")

    # Team and leverage
    s = new_slide(prs)
    title_block(s, "Team & Leverage", "AI-native operating model with focused specialist hiring")
    bullets(
        s,
        [
            "Founder-led vision, product architecture, and quality bar.",
            "AI-assisted development loops compress build-test-iterate cycles.",
            "Next hires: full-stack product engineer, applied AI systems engineer, and field-ops design partner.",
        ],
    )
    footer(s, "Operating principle: high density team, high ownership, high velocity", "07")

    # Ask
    s = new_slide(prs)
    title_block(s, "The Ask", "Partner support to accelerate from strong foundation to category leader")
    bullets(
        s,
        [
            "Capital and network support to scale design-partner deployments.",
            "Introductions to contractors with complex, multi-team execution workflows.",
            "Strategic partnership to accelerate go-to-market and category positioning.",
        ],
    )
    footer(s, "Contact: michael@futurecontractorsofamerica.com", "08")

    # Closing
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s, dark=True)
    title_block(s, "FCA is building the AI-powered operating system", "for how contractors win work and deliver projects.", dark=True)
    bullets(
        s,
        [
            "Future Contractors of America + Auricrux",
            "www.futurecontractorsofamerica.com",
            "github.com/Future-Contractors-of-America-LLC",
        ],
        y=2.55,
        h=2.6,
        size=25,
        dark=True,
    )

    prs.save(OUTPUT)
    return OUTPUT


if __name__ == "__main__":
    out = build()
    print(f"Created: {out}")
