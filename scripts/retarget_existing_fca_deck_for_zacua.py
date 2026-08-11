from pathlib import Path

from pptx import Presentation


SRC = Path(r"C:\Users\Auricrux\OneDrive - Future Contractors of America LLC\Copilot\Created\FCA_NVIDIA_Inception_Pitch_Deck_UNDER_5MB 1.pptx")
DST = Path(r"C:\repos\auricrux-app\FCA_Zacua_Pitch_Deck_Aligned_Final.pptx")


REPLACEMENTS = [
    ("NVIDIA Inception Pitch Deck", "Zacua Ventures Pitch Deck"),
    ("NVIDIA Inception", "Zacua Ventures"),
    ("NVIDIA", "Zacua Ventures"),
]


def replace_in_text_frame(tf):
    # Replace inside runs to preserve most formatting.
    for p in tf.paragraphs:
        for run in p.runs:
            if not run.text:
                continue
            text = run.text
            for old, new in REPLACEMENTS:
                text = text.replace(old, new)
            run.text = text


def replace_in_shape(shape):
    if hasattr(shape, "text_frame") and shape.text_frame is not None:
        replace_in_text_frame(shape.text_frame)
    if shape.has_table:
        for row in shape.table.rows:
            for cell in row.cells:
                replace_in_text_frame(cell.text_frame)
    if shape.shape_type == 6 and hasattr(shape, "shapes"):
        for sub in shape.shapes:
            replace_in_shape(sub)


def main():
    if not SRC.exists():
        raise FileNotFoundError(f"Source deck not found: {SRC}")

    prs = Presentation(str(SRC))

    for slide in prs.slides:
        for shape in slide.shapes:
            replace_in_shape(shape)

    prs.save(str(DST))
    print(f"Created: {DST}")


if __name__ == "__main__":
    main()
