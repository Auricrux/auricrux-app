from pathlib import Path
from pptx import Presentation

SRC = Path(r"C:\repos\auricrux-app\_source_FCA_Lender_Packet_v2_ProductProof.pptx")
DST = Path(r"C:\repos\auricrux-app\FCA_Zacua_Pitch_Deck_Aligned_v5.pptx")

REPLACEMENTS = [
    ("Lender + Backer Packet", "Investor Pitch Deck"),
    ("Lender Packet", "Investor Deck"),
    ("Lender", "Investor"),
    ("Backer", "Partner"),
]


def replace_text(text: str) -> str:
    out = text
    for old, new in REPLACEMENTS:
        out = out.replace(old, new)
    return out


def process_shape(shape):
    if hasattr(shape, "text_frame") and shape.text_frame is not None:
        for p in shape.text_frame.paragraphs:
            for run in p.runs:
                if run.text:
                    run.text = replace_text(run.text)

    if getattr(shape, "has_table", False):
        for row in shape.table.rows:
            for cell in row.cells:
                for p in cell.text_frame.paragraphs:
                    for run in p.runs:
                        if run.text:
                            run.text = replace_text(run.text)

    # Group shapes
    if getattr(shape, "shape_type", None) == 6 and hasattr(shape, "shapes"):
        for sub in shape.shapes:
            process_shape(sub)


def main() -> None:
    if not SRC.exists():
        raise FileNotFoundError(f"Missing source deck: {SRC}")

    prs = Presentation(str(SRC))
    for slide in prs.slides:
        for shape in slide.shapes:
            process_shape(shape)

    prs.save(str(DST))
    print(f"Created: {DST}")


if __name__ == "__main__":
    main()
