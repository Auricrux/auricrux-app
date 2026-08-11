from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


OUTPUT = Path("C:/repos/auricrux-app/FCA_Auricrux_Zacua_Deck_v4.pptx")
BADGE = Path("C:/repos/auricrux-app/apps/auricrux-mobile/assets/auricrux-badge.png")

# Palette
NAVY = RGBColor(10, 24, 51)
BLUE = RGBColor(20, 95, 170)
TEAL = RGBColor(12, 145, 130)
WHITE = RGBColor(255, 255, 255)
BG = RGBColor(246, 249, 253)
TEXT = RGBColor(28, 34, 44)
MUTED = RGBColor(94, 107, 127)


def style_run(run, size=18, bold=False, color=TEXT):
    run.font.name = "Aptos"
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color



def add_bg(slide, dark=False):
    rect = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        Inches(0),
        Inches(0),
        Inches(13.33),
        Inches(7.5),
    )
    rect.fill.solid()
    rect.fill.fore_color.rgb = NAVY if dark else BG
    rect.line.fill.background()

    top = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        Inches(0),
        Inches(0),
        Inches(13.33),
        Inches(0.24),
    )
    top.fill.solid()
    top.fill.fore_color.rgb = TEAL if dark else NAVY
    top.line.fill.background()



def title_block(slide, title, subtitle=None, dark=False):
    t = slide.shapes.add_textbox(Inches(0.78), Inches(0.42), Inches(11.8), Inches(1.4)).text_frame
    t.clear()

    p1 = t.paragraphs[0]
    r1 = p1.add_run()
    r1.text = title
    style_run(r1, size=36, bold=True, color=WHITE if dark else NAVY)

    if subtitle:
        p2 = t.add_paragraph()
        p2.space_before = Pt(6)
        r2 = p2.add_run()
        r2.text = subtitle
        style_run(r2, size=15, color=RGBColor(203, 221, 242) if dark else MUTED)



def bullets(slide, items, x=0.9, y=1.9, w=11.6, h=4.9, size=23, dark=False):
    tf = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h)).text_frame
    tf.clear()
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.space_after = Pt(11)
        p.font.name = "Aptos"
        p.font.size = Pt(size)
        p.font.color.rgb = WHITE if dark else TEXT



def footer(slide, left, right=""):
    l = slide.shapes.add_textbox(Inches(0.78), Inches(7.02), Inches(8.8), Inches(0.3)).text_frame
    l.clear()
    lp = l.paragraphs[0]
    lr = lp.add_run()
    lr.text = left
    style_run(lr, size=10, color=BLUE)

    r = slide.shapes.add_textbox(Inches(9.3), Inches(7.02), Inches(3.2), Inches(0.3)).text_frame
    r.clear()
    rp = r.paragraphs[0]
    rp.alignment = PP_ALIGN.RIGHT
    rr = rp.add_run()
    rr.text = right
    style_run(rr, size=10, color=MUTED)



def card(slide, x, y, w, h, title, value, body, accent):
    shape = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(x),
        Inches(y),
        Inches(w),
        Inches(h),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = WHITE
    shape.line.color.rgb = accent
    shape.line.width = Pt(1.5)

    tf = slide.shapes.add_textbox(Inches(x + 0.2), Inches(y + 0.15), Inches(w - 0.4), Inches(h - 0.3)).text_frame
    tf.clear()

    p0 = tf.paragraphs[0]
    r0 = p0.add_run()
    r0.text = title
    style_run(r0, size=12, bold=True, color=accent)

    p1 = tf.add_paragraph()
    p1.space_before = Pt(5)
    r1 = p1.add_run()
    r1.text = value
    style_run(r1, size=26, bold=True, color=NAVY)

    p2 = tf.add_paragraph()
    p2.space_before = Pt(5)
    r2 = p2.add_run()
    r2.text = body
    style_run(r2, size=11, color=MUTED)



def new_slide(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s)
    return s



def build() -> Path:
    prs = Presentation()

    # 1 Cover
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s, dark=True)
    title_block(
        s,
        "Future Contractors of America",
        "Investor Deck v4 | AI-Powered Construction Operating System",
        dark=True,
    )

    if BADGE.exists():
        s.shapes.add_picture(str(BADGE), Inches(9.45), Inches(1.0), height=Inches(1.8))

    bullets(
        s,
        [
            "One platform connecting leads, bids, projects, field operations, compliance, billing, and workforce readiness.",
            "Auricrux is the embedded intelligence layer that guides what to do next in context.",
        ],
        y=2.4,
        h=2.6,
        size=22,
        dark=True,
    )
    footer(s, "www.futurecontractorsofamerica.com", "v4")

    # 2 Problem
    s = new_slide(prs)
    title_block(s, "The Problem", "Contractor execution is fragmented across tools and teams")
    bullets(
        s,
        [
            "Disconnected workflows force teams to manage work through spreadsheets, inboxes, and manual follow-up.",
            "Lead -> bid -> project transitions lose context, creating delays, rework, and cash-flow drag.",
            "Critical expertise is trapped in senior operators instead of encoded into scalable systems.",
        ],
    )
    footer(s, "Pain point: continuity failure across handoffs", "01")

    # 3 Solution
    s = new_slide(prs)
    title_block(s, "The Solution", "FCA + Auricrux as the system of execution")
    bullets(
        s,
        [
            "FCA unifies daily contractor operations from intake through closeout in a single operational system.",
            "Auricrux continuously interprets workflow signals and recommends next-best actions.",
            "Result: better schedule reliability, cleaner coordination, and stronger billing confidence.",
        ],
    )
    footer(s, "Product lens: operational command, not point-solution sprawl", "02")

    # 4 Why now + proof cards
    s = new_slide(prs)
    title_block(s, "Why Now", "Public product signals show an immediate wedge")

    card(s, 0.9, 1.95, 3.95, 2.05, "Live ROI Framing", "18-32%", "Financial leakage reduction shown in live FCA experience", BLUE)
    card(s, 5.0, 1.95, 3.95, 2.05, "Admin Friction", "$29,825/mo", "Example monthly friction cost in capacity-gap calculator", TEAL)
    card(s, 9.1, 1.95, 3.25, 2.05, "Projected Reclaim", "$10,141/mo", "Example reclaim estimate shown in live workflow narrative", RGBColor(132, 92, 180))

    bullets(
        s,
        [
            "FCA public shell already demonstrates platform depth across bids, projects, finance, academy, and trust routes.",
            "Auricrux is positioned as embedded execution intelligence rather than a standalone chatbot.",
        ],
        y=4.45,
        h=2.0,
        size=20,
    )
    footer(s, "Source: futurecontractorsofamerica.com public pages (Jul 2026 snapshot)", "03")

    # 5 Product depth
    s = new_slide(prs)
    title_block(s, "Product Depth", "Execution modules already mapped across the contractor lifecycle")
    bullets(
        s,
        [
            "Qualification, pipeline, estimating, project flow, field tasks, scheduling, plans, design, and immersive review.",
            "Finance, billing, legal, customer portal, support, and warranty/referral continuity built into the same shell.",
            "Academy/training pathways tied directly to operations instead of separate LMS silos.",
        ],
    )
    footer(s, "Product routes visible in public shell + authenticated workspace", "04")

    # 6 Traction execution evidence
    s = new_slide(prs)
    title_block(s, "Execution Evidence", "Founder velocity and build cadence")

    card(s, 0.9, 1.95, 3.95, 1.9, "GitHub Contributions", "1,056", "Public profile contributions in last year", BLUE)
    card(s, 5.0, 1.95, 3.95, 1.9, "Jul 2026 Commits", "60", "Across 5 repositories", TEAL)
    card(s, 9.1, 1.95, 3.25, 1.9, "PR Record", "33 merged", "In one active repository lane", RGBColor(132, 92, 180))

    bullets(
        s,
        [
            "This is a shipped system with visible execution history, not speculative concept work.",
            "Convergence path centers on FCA Ecosystem as unified universal repository direction.",
        ],
        y=4.25,
        h=2.2,
        size=21,
    )
    footer(s, "Source: github.com/Auricrux and github.com/Future-Contractors-of-America-LLC (Jul 2026 snapshot)", "05")

    # 7 GTM
    s = new_slide(prs)
    title_block(s, "Go-To-Market", "Land where continuity failures are most expensive")
    bullets(
        s,
        [
            "Wedge: bid-to-project transition, document control, and communication integrity for contractor teams.",
            "Expand: finance readiness, customer transparency, compliance, and training continuity.",
            "Retention: become the daily operating surface for execution, not an occasional reporting tool.",
            "Revenue model: subscription by team/workflow depth with module-based expansion.",
        ],
    )
    footer(s, "GTM: pain-first adoption -> account expansion -> system-of-record behavior", "06")

    # 8 Team model
    s = new_slide(prs)
    title_block(s, "Team & Operating Model", "High-leverage, AI-native execution with focused hires")
    bullets(
        s,
        [
            "Founder-led product and systems direction with AI-assisted execution loops.",
            "Near-term hires: product full-stack engineer, applied AI systems engineer, field-ops design partner.",
            "Talent density bar: ownership, speed, rigor, and zero handoff theater.",
        ],
    )
    footer(s, "Small team, high output, compounding execution quality", "07")

    # 9 Ask
    s = new_slide(prs)
    title_block(s, "The Ask", "Partner support to accelerate category leadership")
    bullets(
        s,
        [
            "Capital and partner support to scale design-partner deployments.",
            "Introductions to contractors with complex multi-team workflows and urgent continuity pain.",
            "Strategic support to move from strong product foundation to category-defining platform company.",
        ],
    )
    footer(s, "Contact: michael@futurecontractorsofamerica.com", "08")

    # 10 Closing
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s, dark=True)
    title_block(
        s,
        "FCA is building the AI-powered operating system",
        "for how contractors win work and deliver projects.",
        dark=True,
    )
    bullets(
        s,
        [
            "Future Contractors of America + Auricrux",
            "www.futurecontractorsofamerica.com",
            "github.com/Future-Contractors-of-America-LLC",
        ],
        y=2.5,
        h=2.8,
        size=25,
        dark=True,
    )

    prs.save(OUTPUT)
    return OUTPUT


if __name__ == "__main__":
    out = build()
    print(f"Created: {out}")
