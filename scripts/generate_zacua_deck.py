from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


OUTPUT = Path("C:/repos/auricrux-app/FCA_Auricrux_Zacua_Deck_v2.pptx")

NAVY = RGBColor(12, 24, 48)
STEEL = RGBColor(36, 89, 125)
LIGHT = RGBColor(247, 249, 252)
DARK = RGBColor(28, 34, 44)
MUTED = RGBColor(99, 109, 126)


def set_font(run, size=20, bold=False, color=DARK, name="Aptos"):
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = name


def add_bg(slide):
    bg = slide.shapes.add_shape(
        1,
        Inches(0),
        Inches(0),
        Inches(13.33),
        Inches(7.5),
    )
    bg.fill.solid()
    bg.fill.fore_color.rgb = LIGHT
    bg.line.fill.background()

    top = slide.shapes.add_shape(
        1,
        Inches(0),
        Inches(0),
        Inches(13.33),
        Inches(0.25),
    )
    top.fill.solid()
    top.fill.fore_color.rgb = NAVY
    top.line.fill.background()


def add_title_block(slide, title, subtitle=""):
    tbox = slide.shapes.add_textbox(Inches(0.8), Inches(0.6), Inches(11.7), Inches(1.0))
    tf = tbox.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = title
    set_font(r, size=36, bold=True, color=NAVY)

    if subtitle:
        p2 = tf.add_paragraph()
        r2 = p2.add_run()
        r2.text = subtitle
        set_font(r2, size=16, color=MUTED)


def add_bullets(slide, bullets):
    box = slide.shapes.add_textbox(Inches(0.9), Inches(1.9), Inches(11.8), Inches(4.8))
    tf = box.text_frame
    tf.clear()
    for i, txt in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = txt
        p.level = 0
        p.space_after = Pt(14)
        p.font.size = Pt(24)
        p.font.color.rgb = DARK
        p.font.name = "Aptos"


def add_footer(slide, text):
    box = slide.shapes.add_textbox(Inches(0.9), Inches(7.0), Inches(12), Inches(0.35))
    tf = box.text_frame
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = PP_ALIGN.LEFT
    p.font.size = Pt(11)
    p.font.color.rgb = STEEL
    p.font.name = "Aptos"


def add_slide(prs, title, subtitle, bullets, footer=""):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_title_block(slide, title, subtitle)
    add_bullets(slide, bullets)
    if footer:
        add_footer(slide, footer)


def build():
    prs = Presentation()

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_title_block(
        slide,
        "Future Contractors of America",
        "AI-Powered Construction Operating System | Founder: Michael J. Bartholomew",
    )
    hero = slide.shapes.add_textbox(Inches(0.9), Inches(2.3), Inches(11.7), Inches(2.7))
    tf = hero.text_frame
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = "One platform for leads, bids, projects, field execution, compliance, billing, and workforce readiness."
    set_font(r, size=28, bold=True, color=DARK)
    p2 = tf.add_paragraph()
    r2 = p2.add_run()
    r2.text = "Auricrux is the embedded intelligence layer guiding next-best actions across every workflow handoff."
    set_font(r2, size=20, color=MUTED)
    add_footer(slide, "www.futurecontractorsofamerica.com")

    add_slide(
        prs,
        "Problem",
        "Construction operations are fragmented and high-friction",
        [
            "Critical workflows are split across disconnected tools, inboxes, and spreadsheets.",
            "Lead-to-bid-to-project handoffs break continuity and create avoidable rework.",
            "Institutional know-how is trapped in people rather than operational systems.",
        ],
    )

    add_slide(
        prs,
        "Solution",
        "FCA + Auricrux unifies operations and intelligence",
        [
            "A single construction operating system replacing point-solution sprawl.",
            "Connected workflows across estimating, projects, files, finance, customer comms, and training.",
            "AI-guided execution that turns operational data into real-time decisions.",
        ],
    )

    add_slide(
        prs,
        "Product",
        "Live platform architecture, not slideware",
        [
            "Public product shell + authenticated portal workflows are already live.",
            "Web and mobile operational surfaces support office-field continuity.",
            "FCA Ecosystem is the convergence path to replace legacy fragmented repos.",
        ],
        "Product: https://www.futurecontractorsofamerica.com",
    )

    add_slide(
        prs,
        "Why Now",
        "AI-native operations are becoming mandatory, not optional",
        [
            "Construction still has major workflow fragmentation and low software continuity.",
            "Modern AI systems can now provide reliable in-context operational guidance.",
            "Margin, compliance, and schedule pressure create immediate adoption urgency.",
        ],
    )

    add_slide(
        prs,
        "Go-To-Market",
        "Land in painful workflow bottlenecks, then expand",
        [
            "Start where continuity failures are most expensive: bid-to-project transition.",
            "Win by improving execution confidence in docs, communication, and billing readiness.",
            "Expand account depth via Academy/training and cross-module adoption.",
        ],
    )

    add_slide(
        prs,
        "Progress",
        "Execution velocity and public build artifacts",
        [
            "Active product iteration with live routes, repos, and deployment workflows.",
            "Operating model proven across web, mobile, and ecosystem convergence work.",
            "Founder-led AI-native execution compresses build-test-iterate cycles.",
        ],
        "GitHub: https://github.com/Future-Contractors-of-America-LLC",
    )

    add_slide(
        prs,
        "Business Model",
        "SaaS expansion from workflow command to system-of-record behavior",
        [
            "Subscription platform priced to team size and operational depth.",
            "Expansion via additional workflow modules and intelligence capabilities.",
            "Long-term retention through embedded daily execution and training continuity.",
        ],
    )

    add_slide(
        prs,
        "Vision",
        "Build the default operating system for contractors",
        [
            "From fragmented software stacks to one continuously learning execution platform.",
            "Capture senior operator judgment and operationalize it for every team member.",
            "Make high-quality contractor operations repeatable, scalable, and auditable.",
        ],
    )

    add_slide(
        prs,
        "The Ask",
        "Partner support to accelerate growth and category leadership",
        [
            "Support pilot-to-scale adoption with high-quality contractor design partners.",
            "Accelerate product depth across compliance, finance, and customer continuity.",
            "Help build a category-defining AI-native company in construction.",
        ],
    )

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_title_block(slide, "Thank You", "Contact: michael@futurecontractorsofamerica.com")
    cta = slide.shapes.add_textbox(Inches(0.9), Inches(2.4), Inches(11.7), Inches(2.0))
    tf = cta.text_frame
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = "FCA is building the AI-powered operating system for construction execution."
    set_font(r, size=30, bold=True, color=NAVY)
    p2 = tf.add_paragraph()
    r2 = p2.add_run()
    r2.text = "Product: www.futurecontractorsofamerica.com | GitHub: Future-Contractors-of-America-LLC"
    set_font(r2, size=16, color=STEEL)

    prs.save(OUTPUT)
    return OUTPUT


if __name__ == "__main__":
    out = build()
    print(f"Created: {out}")
