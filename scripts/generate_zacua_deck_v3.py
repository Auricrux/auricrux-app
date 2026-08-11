from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


OUTPUT = Path("C:/repos/auricrux-app/FCA_Auricrux_Zacua_Deck_v3.pptx")

# Brand palette
NAVY = RGBColor(9, 21, 43)
BLUE = RGBColor(23, 92, 170)
TEAL = RGBColor(16, 143, 129)
SLATE = RGBColor(74, 86, 104)
LIGHT_BG = RGBColor(246, 249, 253)
WHITE = RGBColor(255, 255, 255)
TEXT = RGBColor(24, 31, 42)
MUTED = RGBColor(95, 107, 126)


def set_run_font(run, size, bold=False, color=TEXT):
    run.font.name = "Aptos"
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color



def base_slide(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    bg = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        Inches(0), Inches(0), Inches(13.33), Inches(7.5)
    )
    bg.fill.solid()
    bg.fill.fore_color.rgb = LIGHT_BG
    bg.line.fill.background()

    top = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        Inches(0), Inches(0), Inches(13.33), Inches(0.22)
    )
    top.fill.solid()
    top.fill.fore_color.rgb = NAVY
    top.line.fill.background()

    return slide



def add_title(slide, title, subtitle=None):
    tbox = slide.shapes.add_textbox(Inches(0.75), Inches(0.45), Inches(11.9), Inches(1.2))
    tf = tbox.text_frame
    tf.clear()

    p1 = tf.paragraphs[0]
    r1 = p1.add_run()
    r1.text = title
    set_run_font(r1, 36, True, NAVY)

    if subtitle:
        p2 = tf.add_paragraph()
        p2.space_before = Pt(5)
        r2 = p2.add_run()
        r2.text = subtitle
        set_run_font(r2, 15, False, MUTED)



def add_bullet_block(slide, bullets, x=0.9, y=1.8, w=11.5, h=4.9, size=24):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()

    for i, bullet in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = bullet
        p.level = 0
        p.space_after = Pt(12)
        p.font.name = "Aptos"
        p.font.size = Pt(size)
        p.font.color.rgb = TEXT



def add_footer(slide, left, right=""):
    lbox = slide.shapes.add_textbox(Inches(0.75), Inches(7.05), Inches(8.5), Inches(0.3))
    ltf = lbox.text_frame
    ltf.clear()
    lp = ltf.paragraphs[0]
    lr = lp.add_run()
    lr.text = left
    set_run_font(lr, 10, False, BLUE)

    rbox = slide.shapes.add_textbox(Inches(9.5), Inches(7.05), Inches(3.0), Inches(0.3))
    rtf = rbox.text_frame
    rtf.clear()
    rp = rtf.paragraphs[0]
    rp.alignment = PP_ALIGN.RIGHT
    rr = rp.add_run()
    rr.text = right
    set_run_font(rr, 10, False, SLATE)



def add_metric_card(slide, x, y, w, h, title, value, desc, color):
    card = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(x), Inches(y), Inches(w), Inches(h)
    )
    card.fill.solid()
    card.fill.fore_color.rgb = WHITE
    card.line.color.rgb = color
    card.line.width = Pt(1.5)

    tx = slide.shapes.add_textbox(Inches(x + 0.2), Inches(y + 0.15), Inches(w - 0.4), Inches(h - 0.3))
    tf = tx.text_frame
    tf.clear()

    p0 = tf.paragraphs[0]
    r0 = p0.add_run()
    r0.text = title
    set_run_font(r0, 12, True, color)

    p1 = tf.add_paragraph()
    p1.space_before = Pt(6)
    r1 = p1.add_run()
    r1.text = value
    set_run_font(r1, 26, True, NAVY)

    p2 = tf.add_paragraph()
    p2.space_before = Pt(4)
    r2 = p2.add_run()
    r2.text = desc
    set_run_font(r2, 11, False, MUTED)



def add_section_divider(prs, title, subtitle):
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    bg = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        Inches(0), Inches(0), Inches(13.33), Inches(7.5)
    )
    bg.fill.solid()
    bg.fill.fore_color.rgb = NAVY
    bg.line.fill.background()

    accent = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        Inches(0), Inches(0), Inches(13.33), Inches(0.28)
    )
    accent.fill.solid()
    accent.fill.fore_color.rgb = TEAL
    accent.line.fill.background()

    box = slide.shapes.add_textbox(Inches(0.9), Inches(2.2), Inches(11.6), Inches(2.8))
    tf = box.text_frame
    tf.clear()

    p1 = tf.paragraphs[0]
    r1 = p1.add_run()
    r1.text = title
    set_run_font(r1, 44, True, WHITE)

    p2 = tf.add_paragraph()
    p2.space_before = Pt(10)
    r2 = p2.add_run()
    r2.text = subtitle
    set_run_font(r2, 18, False, RGBColor(197, 216, 242))



def build_v3() -> Path:
    prs = Presentation()

    # 1 Cover
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0), Inches(0), Inches(13.33), Inches(7.5))
    bg.fill.solid()
    bg.fill.fore_color.rgb = NAVY
    bg.line.fill.background()

    glow = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(0.7), Inches(1.4), Inches(12.0), Inches(4.8))
    glow.fill.solid()
    glow.fill.fore_color.rgb = RGBColor(14, 37, 75)
    glow.line.fill.background()

    t = slide.shapes.add_textbox(Inches(1.1), Inches(1.9), Inches(11.0), Inches(1.8)).text_frame
    t.clear()
    p = t.paragraphs[0]
    r = p.add_run()
    r.text = "Future Contractors of America"
    set_run_font(r, 48, True, WHITE)

    p2 = t.add_paragraph()
    p2.space_before = Pt(8)
    r2 = p2.add_run()
    r2.text = "AI-Powered Construction Operating System"
    set_run_font(r2, 24, True, RGBColor(183, 228, 255))

    b = slide.shapes.add_textbox(Inches(1.1), Inches(4.5), Inches(11.0), Inches(1.5)).text_frame
    b.clear()
    bp = b.paragraphs[0]
    br = bp.add_run()
    br.text = "Pitch Deck v3 | Founder: Michael J. Bartholomew"
    set_run_font(br, 16, False, RGBColor(205, 216, 236))

    # 2 Problem
    slide = base_slide(prs)
    add_title(slide, "Problem", "Construction operations are fragmented and expensive to coordinate")
    add_bullet_block(slide, [
        "Leads, bids, projects, documents, and billing are managed across disconnected systems.",
        "Operational handoffs fail silently, causing margin leakage, delays, and avoidable rework.",
        "Institutional know-how lives in senior personnel, not in scalable operating infrastructure.",
    ])
    add_footer(slide, "FCA: replacing fragmentation with workflow continuity", "01")

    # 3 Opportunity
    slide = base_slide(prs)
    add_title(slide, "Opportunity", "A system-of-execution layer for one of the world’s largest industries")
    add_metric_card(slide, 0.9, 1.9, 3.9, 2.05, "Industry Context", "$10T+", "Global construction spend scale", BLUE)
    add_metric_card(slide, 4.95, 1.9, 3.9, 2.05, "Operational Gap", "High", "Fragmentation across day-to-day workflows", TEAL)
    add_metric_card(slide, 9.0, 1.9, 3.4, 2.05, "Why Now", "AI", "Practical guidance and automation layer", RGBColor(140, 84, 188))
    add_bullet_block(slide, [
        "The category winner will be the platform that owns execution continuity, not isolated features.",
        "FCA is designed to become the operating core across office, field, and customer-facing workflows.",
    ], y=4.35, h=2.2, size=21)
    add_footer(slide, "Positioning: system of execution for contractors", "02")

    # 4 Product section divider
    add_section_divider(prs, "Product", "From lead intake to closeout in one operational system")

    # 5 Product architecture
    slide = base_slide(prs)
    add_title(slide, "Product Architecture", "One platform, connected modules, continuous intelligence")

    flow = [
        ("Lead & Intake", 0.9),
        ("Bid & Estimate", 3.3),
        ("Project Ops", 5.7),
        ("Finance & Billing", 8.1),
        ("Training & Workforce", 10.5),
    ]

    for label, x in flow:
        box = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(x), Inches(2.3), Inches(2.1), Inches(1.2))
        box.fill.solid()
        box.fill.fore_color.rgb = WHITE
        box.line.color.rgb = BLUE
        box.line.width = Pt(1.2)

        tf = slide.shapes.add_textbox(Inches(x + 0.08), Inches(2.63), Inches(1.95), Inches(0.8)).text_frame
        tf.clear()
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = label
        set_run_font(r, 12, True, NAVY)

    aur = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(2.5), Inches(4.35), Inches(8.2), Inches(1.55))
    aur.fill.solid()
    aur.fill.fore_color.rgb = RGBColor(236, 246, 255)
    aur.line.color.rgb = TEAL
    aur.line.width = Pt(1.5)

    tfa = slide.shapes.add_textbox(Inches(2.75), Inches(4.65), Inches(7.7), Inches(1.0)).text_frame
    tfa.clear()
    p0 = tfa.paragraphs[0]
    p0.alignment = PP_ALIGN.CENTER
    r0 = p0.add_run()
    r0.text = "Auricrux Intelligence Layer"
    set_run_font(r0, 18, True, TEAL)

    p1 = tfa.add_paragraph()
    p1.alignment = PP_ALIGN.CENTER
    r1 = p1.add_run()
    r1.text = "Guides next-best actions, risk signals, and cross-team continuity"
    set_run_font(r1, 12, False, SLATE)

    add_footer(slide, "Live shell: https://www.futurecontractorsofamerica.com", "03")

    # 6 Product proof
    slide = base_slide(prs)
    add_title(slide, "Product Proof", "Live artifacts that demonstrate execution")
    add_bullet_block(slide, [
        "Public platform shell with workflow-first narrative and portal entry points.",
        "Operational modules across bids, projects, files, scheduling, finance, and academy.",
        "Embedded Auricrux layer presented as execution guidance across the system.",
    ], y=1.9, h=2.6, size=22)

    add_metric_card(slide, 0.9, 4.6, 4.0, 1.7, "Live Product", "futurecontractorsofamerica.com", "Public system shell + routing depth", BLUE)
    add_metric_card(slide, 5.1, 4.6, 4.0, 1.7, "Unified Repo Direction", "fca-ecosystem", "Convergence path replacing fragmented repos", TEAL)
    add_metric_card(slide, 9.3, 4.6, 3.1, 1.7, "Execution Cadence", "Active", "Continuous build + iteration", RGBColor(140, 84, 188))
    add_footer(slide, "GitHub org: Future-Contractors-of-America-LLC", "04")

    # 7 GTM + model
    slide = base_slide(prs)
    add_title(slide, "Go-To-Market", "Land on painful bottlenecks, then expand into system-of-record")
    add_bullet_block(slide, [
        "Land: fix bid-to-project continuity and document/communication gaps.",
        "Expand: activate finance readiness, customer visibility, and workforce training loops.",
        "Retain: become the daily command center for contractor execution.",
        "Revenue: subscription model by team/workflow depth with module expansion over time.",
    ], y=1.8, h=4.6, size=22)
    add_footer(slide, "Wedge -> expansion -> platform standard", "05")

    # 8 Traction section divider
    add_section_divider(prs, "Progress", "Built from real workflow pain, now converging into one platform")

    # 9 Progress
    slide = base_slide(prs)
    add_title(slide, "Progress to Date", "From concept to active product system")
    add_bullet_block(slide, [
        "Working product foundation with live modules and public artifacts.",
        "Cross-repo build history demonstrates implementation depth and operational range.",
        "Convergence strategy now centered on FCA Ecosystem as unified universal repo.",
        "Founder-led AI-native operating model accelerates iteration velocity.",
    ], y=1.8, h=4.6, size=22)
    add_footer(slide, "Core link: https://github.com/Future-Contractors-of-America-LLC/fca-ecosystem", "06")

    # 10 Team
    slide = base_slide(prs)
    add_title(slide, "Team & Hiring Plan", "Small, high-density team focused on execution leverage")
    add_metric_card(slide, 0.9, 1.95, 3.9, 1.95, "Role 1", "Full-Stack Product Engineer", "Owns workflow execution end-to-end", BLUE)
    add_metric_card(slide, 5.0, 1.95, 3.9, 1.95, "Role 2", "Applied AI Systems Engineer", "Production-grade agent + eval stack", TEAL)
    add_metric_card(slide, 9.1, 1.95, 3.3, 1.95, "Role 3", "Field Ops Design Partner", "Keeps product anchored to reality", RGBColor(140, 84, 188))

    add_bullet_block(slide, [
        "Talent bar: extreme ownership, high standards, no-silo execution.",
        "Operating model: founder vision + AI leverage + focused specialist hires.",
    ], y=4.35, h=2.0, size=20)
    add_footer(slide, "Hiring for compounding execution quality", "07")

    # 11 Ask
    slide = base_slide(prs)
    add_title(slide, "The Ask", "Partner support to accelerate category leadership")
    add_bullet_block(slide, [
        "Capital and network support to scale design-partner adoption.",
        "Introductions to complex contractor operators with urgent continuity pain.",
        "Strategic support to move from strong product foundation to market-defining platform.",
    ], y=2.0, h=3.8, size=24)
    add_footer(slide, "Contact: michael@futurecontractorsofamerica.com", "08")

    # 12 Closing
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    c_bg = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0), Inches(0), Inches(13.33), Inches(7.5))
    c_bg.fill.solid()
    c_bg.fill.fore_color.rgb = NAVY
    c_bg.line.fill.background()

    c_box = slide.shapes.add_textbox(Inches(0.9), Inches(2.0), Inches(11.6), Inches(3.2)).text_frame
    c_box.clear()
    cp1 = c_box.paragraphs[0]
    cp1.alignment = PP_ALIGN.CENTER
    cr1 = cp1.add_run()
    cr1.text = "FCA is building the AI-powered operating system"
    set_run_font(cr1, 42, True, WHITE)

    cp2 = c_box.add_paragraph()
    cp2.alignment = PP_ALIGN.CENTER
    cp2.space_before = Pt(10)
    cr2 = cp2.add_run()
    cr2.text = "for how contractors win work and deliver projects."
    set_run_font(cr2, 30, True, RGBColor(184, 222, 255))

    cp3 = c_box.add_paragraph()
    cp3.alignment = PP_ALIGN.CENTER
    cp3.space_before = Pt(22)
    cr3 = cp3.add_run()
    cr3.text = "www.futurecontractorsofamerica.com"
    set_run_font(cr3, 17, False, RGBColor(197, 232, 255))

    prs.save(OUTPUT)
    return OUTPUT


if __name__ == "__main__":
    path = build_v3()
    print(f"Created: {path}")
